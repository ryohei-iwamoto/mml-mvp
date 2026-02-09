import argparse
import json
import os
import re

import olefile

try:
    from pptx import Presentation
except Exception:
    Presentation = None


def _extract_pptx(path):
    if Presentation is None:
        raise RuntimeError("python-pptx is not available")
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t:
                    texts.append(t)
        if texts:
            slides.append({"slide": i, "texts": texts})
    return slides


def _extract_ppt(path):
    if not olefile.isOleFile(path):
        raise RuntimeError("Not an OLE (ppt) file")
    ole = olefile.OleFileIO(path)
    if not ole.exists("PowerPoint Document"):
        ole.close()
        raise RuntimeError("PowerPoint Document stream not found")
    data = ole.openstream("PowerPoint Document").read()
    ole.close()

    # ヒューリスティック: UTF-16LEっぽい文字列を一定長以上で抽出。
    texts = []
    try:
        decoded = data.decode("utf-16le", errors="ignore")
        for m in re.finditer(r"[^\x00-\x08\x0b\x0c\x0e-\x1f]{4,}", decoded):
            text = m.group(0).strip()
            if text and len(text) >= 4:
                texts.append(text)
    except Exception:
        texts = []
    # 順序を保ったまま重複を除去。
    seen = set()
    uniq = []
    for t in texts:
        if t in seen:
            continue
        seen.add(t)
        uniq.append(t)
    return [{"slide": None, "texts": uniq}]


def main():
    parser = argparse.ArgumentParser(description="Extract text from .ppt/.pptx")
    parser.add_argument("path", help="Path to .ppt or .pptx")
    parser.add_argument("--json", action="store_true", help="Output JSON")
    args = parser.parse_args()

    path = args.path
    if not os.path.exists(path):
        raise SystemExit(f"File not found: {path}")

    ext = os.path.splitext(path)[1].lower()
    if ext == ".pptx":
        slides = _extract_pptx(path)
    elif ext == ".ppt":
        slides = _extract_ppt(path)
    else:
        raise SystemExit("Only .ppt or .pptx is supported")

    if args.json:
        print(json.dumps(slides, ensure_ascii=False, indent=2))
        return

    for slide in slides:
        if slide["slide"] is None:
            print("PPT Text:")
        else:
            print(f"Slide {slide['slide']}:")
        for t in slide["texts"]:
            print("-", t)
        print()


if __name__ == "__main__":
    main()
